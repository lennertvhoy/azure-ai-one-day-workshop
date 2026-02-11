import io
import os
import uuid
from pathlib import Path
from typing import Any

import requests
from docx import Document
from fastapi import FastAPI, File, UploadFile
from fastapi.responses import HTMLResponse
from openai import AzureOpenAI
from pydantic import BaseModel, Field
from pypdf import PdfReader
from pptx import Presentation

from .search import search_top_k, upload_documents


def get_env(name: str) -> str:
    v = os.getenv(name)
    if not v:
        raise RuntimeError(f"Missing environment variable: {name}")
    return v


def get_aoai_client() -> AzureOpenAI:
    return AzureOpenAI(
        api_key=get_env("AZURE_OPENAI_API_KEY"),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-10-21"),
        azure_endpoint=get_env("AZURE_OPENAI_ENDPOINT"),
    )


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    text = text.strip()
    if not text:
        return []
    out: list[str] = []
    i = 0
    step = max(1, chunk_size - overlap)
    while i < len(text):
        out.append(text[i : i + chunk_size])
        i += step
    return out


def extract_text(filename: str, content: bytes) -> str:
    ext = Path(filename.lower()).suffix
    if ext in {".txt", ".md"}:
        return content.decode("utf-8", errors="ignore")
    if ext == ".pdf":
        reader = PdfReader(io.BytesIO(content))
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    if ext == ".docx":
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)
    if ext == ".pptx":
        prs = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text)
        return "\n".join(lines)
    raise ValueError(f"Unsupported file type: {ext}")


def intake_via_lab1(text: str) -> dict[str, Any]:
    lab1_url = os.getenv("LAB1_URL")
    if not lab1_url:
        return {"doc_type": "unknown", "routing": {"team": "Unknown", "priority": "medium"}}
    r = requests.post(f"{lab1_url.rstrip('/')}/intake", json={"text": text}, timeout=60)
    r.raise_for_status()
    return r.json()


class ChatRequest(BaseModel):
    question: str = Field(min_length=1)


class Citation(BaseModel):
    source: str
    chunk: int | None = None


class ChatResponse(BaseModel):
    answer: str
    citations: list[Citation]
    debug: dict[str, Any] | None = None


SYSTEM_PROMPT = """
You are a policy assistant.
Rules:
- Answer ONLY using the provided SOURCES.
- If the sources do not contain the answer, say: "I don't know based on the provided documents."
- Provide citations in the response as [source#chunk].
""".strip()


app = FastAPI(title="RAG Policy Bot", version="0.2.0")


@app.get("/", response_class=HTMLResponse)
def home() -> str:
    return """
<!doctype html>
<html>
<head>
  <meta charset='utf-8' />
  <title>AI Operations Desk</title>
  <style>
    body { font-family: Arial, sans-serif; margin: 24px; max-width: 980px; }
    textarea, input { width: 100%; margin: 6px 0 12px; }
    button { padding: 8px 12px; margin: 4px 0 12px; }
    pre { background: #f5f5f5; padding: 10px; overflow:auto; }
  </style>
</head>
<body>
  <h1>AI Operations Desk</h1>
  <p>Bulk upload docs into RAG index + ask grounded questions.</p>

  <h2>1) Bulk Upload</h2>
  <input id="files" type="file" multiple />
  <button onclick="upload()">Upload + Index</button>
  <pre id="uploadOut"></pre>

  <h2>2) Ask Policy Bot</h2>
  <input id="question" value="When should I report phishing?" />
  <button onclick="ask()">Ask</button>
  <pre id="chatOut"></pre>

  <script>
    async function upload() {
      const fd = new FormData();
      const files = document.getElementById('files').files;
      for (const f of files) fd.append('files', f);
      const res = await fetch('/upload', { method: 'POST', body: fd });
      document.getElementById('uploadOut').textContent = JSON.stringify(await res.json(), null, 2);
    }
    async function ask() {
      const question = document.getElementById('question').value;
      const res = await fetch('/chat', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({question})
      });
      document.getElementById('chatOut').textContent = JSON.stringify(await res.json(), null, 2);
    }
  </script>
</body>
</html>
"""


@app.get("/health")
def health():
    return {"ok": True}


@app.post("/upload")
async def upload(files: list[UploadFile] = File(...)):
    """Bulk-ingest endpoint used by the web UI.

    Flow per file:
    1) extract raw text from uploaded file
    2) optionally call Lab1 /intake for doc classification
    3) chunk text into searchable pieces
    4) upload chunks to Azure AI Search
    """
    index_name = os.getenv("SEARCH_INDEX", "policy-index")
    all_docs: list[dict[str, Any]] = []
    results: list[dict[str, Any]] = []

    for f in files:
        raw = await f.read()
        try:
            text = extract_text(f.filename or "unknown", raw)
            intake = intake_via_lab1(text)
            doc_type = str(intake.get("doc_type", "unknown"))

            chunks = chunk_text(text)
            for i, ch in enumerate(chunks):
                all_docs.append(
                    {
                        "id": str(uuid.uuid4()),
                        "content": ch,
                        "source": f.filename or "unknown",
                        "title": doc_type,
                        "chunk": i,
                    }
                )
            results.append({"file": f.filename, "status": "ok", "chunks": len(chunks), "doc_type": doc_type})
        except Exception as e:
            results.append({"file": f.filename, "status": "error", "error": str(e)})

    if all_docs:
        upload_documents(index_name=index_name, documents=all_docs)

    return {
        "ok": True,
        "files": len(files),
        "chunks_uploaded": len(all_docs),
        "results": results,
    }


@app.post("/chat", response_model=ChatResponse)
def chat(req: ChatRequest):
    """RAG chat endpoint: retrieve relevant chunks, then generate grounded answer."""
    index_name = os.getenv("SEARCH_INDEX", "policy-index")
    top_k = int(os.getenv("SEARCH_TOP_K", "5"))

    hits = search_top_k(index_name=index_name, query=req.question, top=top_k)

    sources_lines: list[str] = []
    citations: list[Citation] = []

    for h in hits:
        src = str(h.get("source") or "unknown")
        chunk = h.get("chunk")
        content = str(h.get("content") or "")
        citations.append(Citation(source=src, chunk=chunk if isinstance(chunk, int) else None))
        sources_lines.append(f"SOURCE: {src}#{chunk}\n{content}")

    sources_block = "\n\n".join(sources_lines)

    client = get_aoai_client()
    deployment = get_env("AZURE_OPENAI_DEPLOYMENT")

    user_prompt = f"QUESTION:\n{req.question}\n\nSOURCES:\n{sources_block}"

    resp = client.chat.completions.create(
        model=deployment,
        temperature=0.2,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ],
    )

    answer = (resp.choices[0].message.content or "").strip()

    return ChatResponse(answer=answer, citations=citations, debug={"hits": len(hits)})
